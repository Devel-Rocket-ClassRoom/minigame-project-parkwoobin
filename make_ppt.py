from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG        = RGBColor(0x1A, 0x1A, 0x2E)   # 네이비 배경
C_ACCENT    = RGBColor(0xE9, 0x4F, 0x37)   # 레드 오렌지 포인트
C_ACCENT2   = RGBColor(0xF5, 0xA6, 0x23)   # 골드 서브 포인트
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)
C_CARD      = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
C_CARD2     = RGBColor(0x0F, 0x3D, 0x54)   # 카드 배경 2

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 헬퍼 함수들 ─────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    return txBox

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide(prs):
    s = prs.slides.add_slide(blank)
    set_bg(s, C_BG)
    return s

def add_title_bar(slide, title, subtitle=None):
    """상단 타이틀 바"""
    add_rect(slide, 0, 0, 13.33, 1.1, C_ACCENT)
    add_text(slide, title, 0.3, 0.1, 10, 0.7,
             font_size=32, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.75, 12, 0.4,
                 font_size=13, bold=False, color=RGBColor(0xFF,0xEE,0xCC))

def add_card(slide, x, y, w, h, title, items, color=None):
    """정보 카드"""
    bg_color = color or C_CARD
    add_rect(slide, x, y, w, h, bg_color)
    # 카드 상단 강조선
    add_rect(slide, x, y, w, 0.05, C_ACCENT2)
    add_text(slide, title, x+0.15, y+0.08, w-0.3, 0.4,
             font_size=14, bold=True, color=C_ACCENT2)
    for i, item in enumerate(items):
        add_text(slide, f"• {item}", x+0.15, y+0.5+i*0.38, w-0.3, 0.38,
                 font_size=12, color=C_LIGHT)

def add_section_label(slide, text, x, y):
    add_rect(slide, x, y, 0.07, 0.32, C_ACCENT)
    add_text(slide, text, x+0.15, y, 6, 0.35,
             font_size=15, bold=True, color=C_ACCENT2)


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════════════════════════
s = add_slide(prs)

# 배경 장식 사각형
add_rect(s, 0, 0, 13.33, 7.5, C_BG)
add_rect(s, 9.5, 0, 3.83, 7.5, C_CARD)
add_rect(s, 9.5, 0, 0.06, 7.5, C_ACCENT)

# 메인 타이틀
add_text(s, "골목고양이 생존기", 0.6, 1.6, 9, 1.2,
         font_size=52, bold=True, color=C_WHITE)
add_text(s, "Alley Cat Survival", 0.6, 2.8, 8, 0.6,
         font_size=22, bold=False, color=C_ACCENT2)

# 구분선
add_rect(s, 0.6, 3.55, 5.5, 0.06, C_ACCENT)

# 서브 정보
add_text(s, "Unity 2D 횡스크롤 생존 로그라이트", 0.6, 3.75, 8, 0.5,
         font_size=16, color=C_LIGHT)
add_text(s, "Android 우선  |  v0.3.1", 0.6, 4.2, 8, 0.45,
         font_size=14, color=C_LIGHT)

# 우측 패널 정보
add_text(s, "빌드 정보", 9.8, 1.5, 3, 0.45,
         font_size=14, bold=True, color=C_ACCENT2)
infos = [
    "버전   v0.3.1",
    "플랫폼  Android",
    "엔진   Unity 2D",
    "장르   생존 로그라이트",
    "개발   1인 개발",
    "릴리스  2026-06-08",
]
for i, info in enumerate(infos):
    add_text(s, info, 9.8, 2.1+i*0.5, 3.3, 0.45, font_size=12, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 게임 개요
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "게임 개요", "Game Overview")

add_section_label(s, "스토리", 0.4, 1.4)
add_text(s,
    "교통사고로 혼수상태에 빠진 김경일. 눈을 뜨니 고양이의 몸속이었다.\n"
    "골목을 헤쳐나가 자신의 몸으로 돌아가야 한다.",
    0.65, 1.4, 8.5, 0.9, font_size=13, color=C_LIGHT)

add_section_label(s, "진행 구조", 0.4, 2.55)
stages = ["프롤로그", "튜토리얼", "Map1\n도심 골목", "Map2\n야간 골목", "Map3\n도시", "Map4\n하수구", "Map5\n병원\n(엔딩)"]
colors = [C_CARD2, C_CARD2, C_CARD, C_CARD, C_CARD, C_CARD, RGBColor(0x1A,0x3A,0x1A)]
for i, (stage, col) in enumerate(zip(stages, colors)):
    bx = 0.4 + i * 1.84
    add_rect(s, bx, 2.85, 1.6, 1.0, col)
    add_rect(s, bx, 2.85, 1.6, 0.05, C_ACCENT if i < 2 else C_ACCENT2)
    add_text(s, stage, bx+0.05, 2.95, 1.5, 0.85, font_size=11, align=PP_ALIGN.CENTER, color=C_LIGHT)
    if i < len(stages)-1:
        add_text(s, "→", bx+1.62, 3.2, 0.2, 0.5, font_size=16, bold=True, color=C_ACCENT2)

add_section_label(s, "핵심 루프", 0.4, 4.15)
loops = [
    ("이동 & 탐색", "골목을 달리고 점프하며\n적과 아이템을 탐색"),
    ("전투", "공격·대시·턴·벽점프로\n적을 처치"),
    ("생존 관리", "체력·배고픔을 유지하며\n체크포인트 확보"),
    ("성장", "코인으로 상점에서\n스탯 업그레이드"),
]
for i, (title, desc) in enumerate(loops):
    bx = 0.4 + i * 3.1
    add_rect(s, bx, 4.45, 2.8, 1.45, C_CARD)
    add_rect(s, bx, 4.45, 2.8, 0.05, C_ACCENT2)
    add_text(s, f"0{i+1}. {title}", bx+0.15, 4.52, 2.5, 0.38, font_size=13, bold=True, color=C_ACCENT2)
    add_text(s, desc, bx+0.15, 4.92, 2.5, 0.85, font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 플레이어 구현
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "플레이어 구현", "Player Implementation")

skills = [
    ("공격", [
        "공격 버튼으로 히트박스 활성화",
        "전방 범위 판정",
        "공격력 업그레이드 연동",
        "공격 쿨타임 시스템",
    ]),
    ("더블 점프", [
        "지면 점프 후 공중에서 1회 추가 점프",
        "착지 시 자동 초기화",
        "벽 점프 후 더블점프 차단",
        "업그레이드 해금 시스템",
    ]),
    ("벽 점프", [
        "벽 접촉 시 슬라이딩 & 고정",
        "점프 입력으로 벽 반대 방향 도약",
        "방향 입력 잠금으로 관성 유지",
        "계단 오판 방지 Raycast 처리",
    ]),
    ("턴", [
        "순간 뒤돌기 동작",
        "무적 판정 시간 부여",
        "턴 쿨타임 업그레이드 연동",
        "이동 중 / 공중 모두 사용 가능",
    ]),
    ("대시", [
        "전방 고속 이동 (중력 제거)",
        "공중 대시 직선 운동",
        "대시 중 무적 없음",
        "대시 쿨타임 업그레이드 연동",
    ]),
]

cols = 5
card_w = 2.4
card_h = 2.8
for i, (title, items) in enumerate(skills):
    bx = 0.3 + i * (card_w + 0.18)
    add_rect(s, bx, 1.3, card_w, card_h, C_CARD)
    add_rect(s, bx, 1.3, card_w, 0.06, C_ACCENT)
    add_text(s, title, bx+0.15, 1.38, card_w-0.3, 0.42,
             font_size=15, bold=True, color=C_ACCENT)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", bx+0.15, 1.88+j*0.45, card_w-0.25, 0.42,
                 font_size=11, color=C_LIGHT)

# 하단 — 생존 시스템
add_section_label(s, "생존 시스템", 0.3, 4.35)
survival = [
    "체력(HP) — 피격 시 감소, 체크포인트 저장",
    "배고픔 — 시간·행동 시 소모. 0 도달 시 HP 지속 감소",
    "체크포인트 — 깃발 접촉으로 위치·스탯 저장",
    "게임오버 — 마지막 체크포인트에서 재시작",
]
for i, t in enumerate(survival):
    add_text(s, f"• {t}", 0.45+i*3.18, 4.72, 3.05, 0.85,
             font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 적 구현
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "적 구현", "Enemy Implementation")

enemies = [
    ("고양이 적", C_CARD, [
        "좌우 순찰 이동",
        "플레이어 근접 시 공격",
        "피격 후 경직 & 넉백",
        "아이템 드롭",
    ]),
    ("개 적", C_CARD, [
        "플레이어 추적 이동",
        "돌진 공격 패턴",
        "대시로 빠른 접근",
        "아이템 드롭",
    ]),
    ("사람 적", C_CARD, [
        "원거리 투사체 공격",
        "일정 거리 유지 이동",
        "점프로 장애물 회피",
        "아이템 드롭",
    ]),
    ("거미 (NEW)", C_CARD2, [
        "Map4 하수구 전용",
        "이동 패턴",
        "공격 패턴",
        "SFX 일부 미완성",
    ]),
]

card_w = 2.95
for i, (name, col, items) in enumerate(enemies):
    bx = 0.3 + i*(card_w+0.2)
    is_new = "NEW" in name
    add_rect(s, bx, 1.3, card_w, 2.7, col)
    add_rect(s, bx, 1.3, card_w, 0.06, C_ACCENT2 if is_new else C_ACCENT)
    add_text(s, name, bx+0.15, 1.38, card_w-0.3, 0.42,
             font_size=14, bold=True, color=C_ACCENT2 if is_new else C_WHITE)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", bx+0.15, 1.9+j*0.45, card_w-0.25, 0.42,
                 font_size=12, color=C_LIGHT)

# 공통 AI 시스템
add_section_label(s, "공통 AI 시스템", 0.3, 4.25)
ai_items = [
    "지면 감지 — 낙사 방지 Raycast",
    "피격 경직 & 무적 프레임",
    "HP → 사망 애니메이션 → 드롭 → 제거",
    "스포너 — 구역별 자동 스폰 & 최대 수 제한",
]
for i, t in enumerate(ai_items):
    col_ = 0 if i < 2 else 1
    row_ = i % 2
    add_text(s, f"• {t}", 0.5 + col_*6.3, 4.62+row_*0.48, 6.0, 0.45,
             font_size=12, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — 아이템 시스템
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "아이템 시스템", "Item System  —  총 19종 (Box 포함)")

categories = [
    ("회복 아이템", C_CARD, [
        "생선 — 배고픔 소량 회복",
        "닭다리 — 배고픔 중간 회복",
        "캔 음식 — 배고픔 대량 회복",
        "우유 — HP + 배고픔 동시 회복",
        "반창고 — HP 소량 회복",
        "붕대 — HP 중간 회복",
    ]),
    ("코인 & 자원", C_CARD, [
        "동전 — 기본 코인",
        "은화 — 중간 코인",
        "금화 — 고급 코인",
        "열쇠 — Box 개봉에 사용",
    ]),
    ("상자 & 특수", C_CARD2, [
        "Box — 열쇠로 개봉, 내부 아이템 드롭",
        "상점 — 코인으로 스탯 업그레이드",
        "",
        "총 19종 구성",
    ]),
    ("상점 업그레이드", C_CARD, [
        "스피드 (Lv.1~5)",
        "대시 속도 (Lv.1~5)",
        "점프 높이 (Lv.1~5)",
        "대시 쿨타임 (Lv.1~5)",
        "턴 쿨타임 (Lv.1~5)",
        "체력 회복 (소모품)",
        "배고픔 회복 (소모품)",
    ]),
]

card_w = 2.98
for i, (title, col, items) in enumerate(categories):
    bx = 0.25 + i*(card_w+0.15)
    add_rect(s, bx, 1.25, card_w, 5.55, col)
    add_rect(s, bx, 1.25, card_w, 0.06, C_ACCENT2)
    add_text(s, title, bx+0.15, 1.33, card_w-0.3, 0.42,
             font_size=13, bold=True, color=C_ACCENT2)
    for j, item in enumerate(items):
        if item:
            add_text(s, f"• {item}", bx+0.15, 1.85+j*0.55, card_w-0.25, 0.52,
                     font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 맵 & 레벨 디자인
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "맵 & 레벨 디자인", "Map & Level Design")

maps = [
    ("프롤로그", "도입부 컷씬\n조작 설명 없음"),
    ("튜토리얼", "기본 조작\n순차 안내"),
    ("Map1\n도심 골목", "낮 배경\n기본 적 등장"),
    ("Map2\n야간 골목", "밤 배경\n난이도 상승"),
    ("Map3\n도시", "복잡한 지형\n다수 적"),
    ("Map4\n하수구", "어두운 실내\n거미 전용 적"),
    ("Map5\n병원", "최종 스테이지\n엔딩 컷씬"),
]
mcolors = [C_CARD2]*2 + [C_CARD]*5
mcolors[-1] = RGBColor(0x1A,0x3A,0x1A)
card_w2 = 1.72
for i, ((name, desc), col) in enumerate(zip(maps, mcolors)):
    bx = 0.25 + i*(card_w2+0.1)
    add_rect(s, bx, 1.3, card_w2, 1.7, col)
    add_rect(s, bx, 1.3, card_w2, 0.05,
             C_ACCENT if i < 2 else C_ACCENT2)
    add_text(s, name, bx+0.1, 1.37, card_w2-0.2, 0.55,
             font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, bx+0.1, 1.95, card_w2-0.2, 0.85,
             font_size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

# 레벨 디자인 요소
add_section_label(s, "레벨 디자인 요소", 0.3, 3.25)
ld_items = [
    ("화살표 안내", "각 맵 진행 방향을\n화살표로 표시"),
    ("Help 블록", "도움말 팝업 트리거.\n조작법을 인게임에서 안내"),
    ("체크포인트", "깃발 접촉으로\n위치 자동 저장"),
    ("SaveSpot", "저장 위치를 깃발 좌표로\n정확히 기록"),
    ("구역 전환", "구역 끝 빛 오브젝트로\n다음 맵 이동"),
    ("DeadZone", "낙사 구간 감지 후\n게임오버 처리"),
]
for i, (title, desc) in enumerate(ld_items):
    col_ = i % 3
    row_ = i // 3
    bx = 0.3 + col_*4.3
    by = 3.65 + row_*1.65
    add_rect(s, bx, by, 4.0, 1.45, C_CARD)
    add_rect(s, bx, by, 4.0, 0.05, C_ACCENT2)
    add_text(s, title, bx+0.15, by+0.1, 3.7, 0.38,
             font_size=13, bold=True, color=C_ACCENT2)
    add_text(s, desc, bx+0.15, by+0.52, 3.7, 0.78,
             font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — UI & 시스템
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "UI & 시스템", "UI & Core Systems")

systems = [
    ("HUD", [
        "HP 바 & 배고픔 바",
        "코인 & 열쇠 카운터",
        "스킬 쿨타임 아이콘",
        "도움말 버튼",
    ]),
    ("설정", [
        "BGM / SFX 볼륨 슬라이더",
        "언어 전환 (한국어 / English)",
        "메인으로 / 게임 종료",
        "씬 전환 시 설정 유지",
    ]),
    ("게임오버 패널", [
        "사망 후 딜레이 표시",
        "다시하기 — 체크포인트 재시작",
        "메인으로 이동",
        "게임 오버 SFX",
    ]),
    ("저장 시스템", [
        "슬롯 3개 세이브 지원",
        "위치·HP·배고픔·코인·열쇠 저장",
        "체크포인트 좌표 정확 저장",
        "업그레이드 레벨 유지",
    ]),
    ("다국어", [
        "한국어 / English 지원",
        "strings.json 기반 로컬라이징",
        "LocalizedText 컴포넌트 자동 적용",
        "설정 변경 즉시 반영",
    ]),
    ("오디오", [
        "BGM — 씬별 자동 재생",
        "SFX — 행동별 효과음",
        "컷씬 전용 BGM 전환",
        "볼륨 PlayerPrefs 저장",
    ]),
]

card_w3 = 3.95
card_h3 = 2.6
for i, (title, items) in enumerate(systems):
    col_ = i % 3
    row_ = i // 3
    bx = 0.25 + col_*(card_w3+0.22)
    by = 1.25 + row_*(card_h3+0.2)
    add_rect(s, bx, by, card_w3, card_h3, C_CARD)
    add_rect(s, bx, by, card_w3, 0.06, C_ACCENT)
    add_text(s, title, bx+0.15, by+0.1, card_w3-0.3, 0.42,
             font_size=14, bold=True, color=C_ACCENT)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", bx+0.15, by+0.6+j*0.45, card_w3-0.25, 0.42,
                 font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 빌드 현황 & 다음 계획
# ════════════════════════════════════════════════════════════
s = add_slide(prs)
add_title_bar(s, "빌드 현황 & 다음 계획", "Build Status & Roadmap")

# 미구현 / 이슈
add_section_label(s, "미구현 · 이슈", 0.3, 1.3)
issues = [
    ("거미 SFX 미완성", "효과음 일부 미적용 상태", "🔴"),
    ("밸런싱 미조정", "맵 추가 후 난이도 곡선 재설계 필요", "🔴"),
    ("벽 점프 물리", "점프 후 수직 낙하 — 자연스러운 밀려남 미구현", "🟡"),
    ("레벨 디자인 검수", "이동 가능 구역 체크리스트 미정립", "🟡"),
]
for i, (title, desc, icon) in enumerate(issues):
    bx = 0.3 + (i%2)*6.3
    by = 1.68 + (i//2)*0.88
    add_rect(s, bx, by, 6.1, 0.75, C_CARD)
    add_text(s, f"{icon} {title}", bx+0.15, by+0.05, 5.8, 0.35,
             font_size=13, bold=True, color=C_WHITE)
    add_text(s, desc, bx+0.15, by+0.4, 5.8, 0.3,
             font_size=11, color=C_LIGHT)

# v0.3.2 계획
add_section_label(s, "v0.3.2 다음 빌드 계획", 0.3, 3.6)
plans = [
    ("벽 점프 개선", "벽 반대 방향으로\n자연스럽게 밀려나도록"),
    ("수치 밸런싱", "아이템·코인·점프력\n전반 재조정"),
    ("거미 SFX", "미완성 효과음\n마무리"),
    ("통합 테스트", "Map1~5 전 구간\n이동·전투 검증"),
]
for i, (title, desc) in enumerate(plans):
    bx = 0.3 + i*3.15
    add_rect(s, bx, 3.95, 2.9, 1.75, C_CARD2)
    add_rect(s, bx, 3.95, 2.9, 0.05, C_ACCENT2)
    add_text(s, f"0{i+1}. {title}", bx+0.15, 4.02, 2.6, 0.4,
             font_size=13, bold=True, color=C_ACCENT2)
    add_text(s, desc, bx+0.15, 4.48, 2.6, 0.95,
             font_size=11, color=C_LIGHT)

# 일정
add_text(s, "📅  일정: 1주일  (Day1~2 벽점프 → Day3 테스트 → Day4~5 밸런싱 → Day6 SFX → Day7 빌드)",
         0.3, 5.88, 12.5, 0.45, font_size=12, color=C_ACCENT2)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
output = r"D:\Unity_Project\Cat_Project\골목고양이_빌드보고서_v0.3.1.pptx"
prs.save(output)
print("저장 완료: " + output)
